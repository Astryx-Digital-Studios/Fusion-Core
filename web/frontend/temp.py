from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# --- Slide 1: Iridescent Beetle Inspiration ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Mood Board: Iridescent Beetle"

# Placeholder rectangles
for i, label in enumerate(["Beetle Shell Image", "Close-up Iridescence", "Jewellery Inspiration"]):
    left = Inches(0.5 + (i * 3))
    top = Inches(1.5)
    width = Inches(2.5)
    height = Inches(2)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    textbox = slide.shapes.add_textbox(left, top + Inches(2.1), width, Inches(0.5))
    textbox.text = label

# --- Slide 2: Dragonfly Wing Structure ---
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
slide2.shapes.title.text = "Mood Board: Dragonfly Wings"

for i, label in enumerate(["Wing Structure Image", "Vein Pattern Close-up", "Filigree Jewellery Reference"]):
    left = Inches(0.5 + (i * 3))
    top = Inches(1.5)
    width = Inches(2.5)
    height = Inches(2)
    shape = slide2.shapes.add_shape(1, left, top, width, height)
    textbox = slide2.shapes.add_textbox(left, top + Inches(2.1), width, Inches(0.5))
    textbox.text = label

# --- Slide 3: Ant Segmentation + Design Ideas ---
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
slide3.shapes.title.text = "Mood Board: Ant Segmentation"

for i, label in enumerate(["Ant Body Image", "Segment Detail", "Necklace Concept Sketch"]):
    left = Inches(0.5 + (i * 3))
    top = Inches(1.5)
    width = Inches(2.5)
    height = Inches(2)
    shape = slide3.shapes.add_shape(1, left, top, width, height)
    textbox = slide3.shapes.add_textbox(left, top + Inches(2.1), width, Inches(0.5))
    textbox.text = label

# Save file
file_path = "./Insect_Moodboard.pptx"
prs.save(file_path)

file_path
